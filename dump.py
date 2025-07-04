import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE

def safe_deep_dump(obj, max_depth=2, _visited=None, _depth=0):
    if _visited is None:
        _visited = set()

    if id(obj) in _visited:
        return "[Circular Reference]"
    _visited.add(id(obj))

    if isinstance(obj, (str, int, float, bool, type(None))):
        return obj

    if _depth >= max_depth:
        return f"<{type(obj).__name__}>"

    result = {}
    for attr in dir(obj):
        if attr.startswith("_"):
            continue
        try:
            value = getattr(obj, attr)
            if callable(value):
                continue
            result[attr] = safe_deep_dump(value, max_depth, _visited, _depth + 1)
        except Exception as e:
            result[attr] = f"[Error: {e}]"
    return result

def get_rgb_safe(color_obj):
    """Trích màu RGB an toàn từ color object (font.color hoặc fill.fore_color)."""
    if color_obj is None:
        return "None"

    try:
        if color_obj.rgb:  # Đây là kiểu RGBColor → an toàn
            return str(color_obj.rgb)
    except AttributeError:
        pass

    try:
        if color_obj.type == 2 and color_obj.theme_color:
            return f"Theme:{color_obj.theme_color.name}"
    except:
        pass

    return "Theme color or undefined"

def extract_slide_data(pptx_path, for_txt=False):
    prs = Presentation(pptx_path)
    slides = []

    for i, slide in enumerate(prs.slides):
        slide_info = {
            "slide_number": i + 1,
            "shapes": []
        }

        for j, shape in enumerate(slide.shapes):
            shape_info = {
                "shape_index": j + 1,
                "type": shape.shape_type,
                "position": {
                    "x": shape.left,
                    "y": shape.top,
                    "width": shape.width,
                    "height": shape.height
                },
                "background_fill_color": None,
                "border": {},
                "text": [],
                "table": None
            }
            shape_info["raw_attributes"] = safe_deep_dump(shape, max_depth=2)
            # Fill color
            try:
                if shape.fill and shape.fill.fore_color:
                    shape_info["background_fill_color"] = get_rgb_safe(shape.fill.fore_color)
            except:
                shape_info["background_fill_color"] = "Error reading"

            # Border
            try:
                if shape.line:
                    line = shape.line
                    shape_info["border"] = {
                        "color": get_rgb_safe(line.color),
                        "width_pt": round(line.width / 12700, 2) if line.width else "Default",
                        "style": str(line.dash_style) if line.dash_style else "None"
                    }
            except:
                shape_info["border"] = {"error": "Error reading"}

            # Text
            if shape.has_text_frame:
                tf = shape.text_frame
                for p_idx, para in enumerate(tf.paragraphs):
                    para_info = {
                        "paragraph_index": p_idx + 1,
                        "text": para.text.strip().replace("\n", "\\n") if for_txt else para.text.strip(),
                        "runs": []
                    }
                    for run_idx, run in enumerate(para.runs):
                        font = run.font
                        run_info = {
                            "run_index": run_idx + 1,
                            "text": run.text.replace("\n", "\\n") if for_txt else run.text,
                            "font_name": font.name,
                            "font_size": font.size.pt if font.size else None,
                            "bold": font.bold,
                            "italic": font.italic,
                            "font_color": get_rgb_safe(font.color)
                        }
                        para_info["runs"].append(run_info)
                    shape_info["text"].append(para_info)

            # Table
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                try:
                    tbl = shape.table
                    table_data = []
                    for row in tbl.rows:
                        row_data = [cell.text.strip().replace("\n", "\\n") if for_txt else cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    shape_info["table"] = {
                        "rows": len(tbl.rows),
                        "cols": len(tbl.columns),
                        "data": table_data
                    }
                except Exception as e:
                    shape_info["table"] = {"error": str(e)}

            slide_info["shapes"].append(shape_info)
        slides.append(slide_info)
    result = {
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": slides
    }
    return result


# ✅ Dump ra TXT dễ đọc
def describe_pptx_to_txt(pptx_path, output_txt):
    slides_data = extract_slide_data(pptx_path, for_txt=True)
    slides = slides_data["slides"]
    lines = [f"📊 Tổng số slide: {len(slides)}\n"]

    for slide in slides:
        lines.append(f"--- Slide {slide['slide_number']} ---")
        for shape in slide["shapes"]:
            lines.append(f"  📌 Shape {shape['shape_index']}:")
            lines.append(f"    - Type: {shape['type']}")
            pos = shape["position"]
            lines.append(f"    - Position: (x={pos['x']}, y={pos['y']}), size=(w={pos['width']}, h={pos['height']})")
            lines.append(f"    - Background fill color: {shape['background_fill_color']}")

            border = shape.get("border", {})
            if "error" in border:
                lines.append("    - Border info: [error reading]")
            else:
                lines.append(f"    - Border color: {border['color']}")
                lines.append(f"    - Border width: {border['width_pt']} pt")
                lines.append(f"    - Border style: {border['style']}")

            # Text
            for para in shape["text"]:
                lines.append(f"    - Paragraph {para['paragraph_index']}: \"{para['text']}\"")
                for run in para["runs"]:
                    lines.append(f"      ▸ Run {run['run_index']}: \"{run['text']}\"")
                    lines.append(f"        - Font name: {run['font_name']}")
                    lines.append(f"        - Font size: {run['font_size']}")
                    lines.append(f"        - Bold: {run['bold']}")
                    lines.append(f"        - Italic: {run['italic']}")
                    lines.append(f"        - Font color: {run['font_color']}")

            # Table
            table = shape["table"]
            if table:
                if "error" in table:
                    lines.append(f"    - Table content: [error reading] ({table['error']})")
                else:
                    lines.append(f"    - Table content ({table['rows']} rows x {table['cols']} cols):")
                    for r_idx, row in enumerate(table["data"]):
                        row_text = [f"[{c_idx+1}] {val}" for c_idx, val in enumerate(row)]
                        lines.append(f"      Row {r_idx+1}: " + " | ".join(row_text))
        lines.append("")

    with open(output_txt, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))
    print(f"✅ Đã ghi log TXT vào: {output_txt}")

# ✅ Dump ra JSON
def describe_pptx_to_json(pptx_path, output_json):
    data = extract_slide_data(pptx_path, for_txt=False)
    with open(output_json, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)
    print(f"✅ Đã ghi JSON vào: {output_json}")

# Ví dụ sử dụng
if __name__ == "__main__":
    #describe_pptx_to_txt("template\\Pre_DOI_Form_05_2024.pptx", "bin\\dump_output.txt")
    describe_pptx_to_json("template\\Pre_DOI_Form_05_2024.pptx", "bin\\dump_output.json")