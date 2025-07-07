# Refactored build.py with modular table and text rebuild logic

import json
import os
from pptx import Presentation
from pptx.slide import Slide
from pptx.shapes.picture import Picture
from pptx.table import _Cell
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.shapes.autoshape import Shape
from pptx.text.text import _Run
from pptx.text.text import _Paragraph
from pptx.text.text import TextFrame
from data.pptxdata import *
from dacite import from_dict

EMU = 1  # đơn vị đã là EMU trong JSON dump


def apply_cell_border(cell: _Cell, border_info: DL_CellBorder):
    tcPr = cell._tc.get_or_add_tcPr()
    side_map = {
        "left": ("a:lnL", border_info.left),
        "right": ("a:lnR", border_info.right),
        "top": ("a:lnT", border_info.top),
        "bottom": ("a:lnB", border_info.bottom),
        "tl2br": ("a:lnTlToBr", border_info.diagonal_down),
        "bl2tr": ("a:lnBlToTr", border_info.diagonal_up),
    }

    for tag, side_border in side_map.values():
        if side_border is None or side_border.color == 'None':
            continue

        ln = OxmlElement(tag)

        # W + style mặc định giống origin
        width = side_border.width
        if isinstance(width, (int, float)):
            ln.set("w", str(int(width * 12700)))
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        # Color
        color_info = parse_color(side_border.color)
        if color_info["type"] == "rgb":
            solidFill = OxmlElement("a:solidFill")
            srgbClr = OxmlElement("a:srgbClr")
            rgb_val = color_info["value"]
            srgbClr.set(
                "val", f"{rgb_val[0]:02X}{rgb_val[1]:02X}{rgb_val[2]:02X}")
            solidFill.append(srgbClr)
            ln.append(solidFill)
        elif color_info["type"] == "theme":
            solidFill = OxmlElement("a:solidFill")
            schemeClr = OxmlElement("a:schemeClr")
            schemeClr.set("val", color_info["value"].name.lower())
            solidFill.append(schemeClr)
            ln.append(solidFill)

        # Dash style (mặc định solid)
        dash_style = (side_border.dash_type or "solid").lower()
        prstDash = OxmlElement("a:prstDash")
        prstDash.set("val", dash_style)
        ln.append(prstDash)

        # round (origin có)
        ln.append(OxmlElement("a:round"))

        # headEnd và tailEnd (mặc định none)
        headEnd = OxmlElement("a:headEnd")
        headEnd.set("type", "none")
        headEnd.set("w", "med")
        headEnd.set("len", "med")
        ln.append(headEnd)

        tailEnd = OxmlElement("a:tailEnd")
        tailEnd.set("type", "none")
        tailEnd.set("w", "med")
        tailEnd.set("len", "med")
        ln.append(tailEnd)

        tcPr.append(ln)
    return tcPr


def parse_color(color_str: Optional[str]) -> dict:
    if color_str is None or color_str == "None":
        return {"type": "none"}
    if color_str.startswith("RGB:"):
        hex_str = color_str[4:].strip()
        return {
            "type": "rgb",
            "value": RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        }
    if color_str.startswith("Theme:"):
        theme_name = color_str[6:].strip()
        return {
            "type": "theme",
            "value": getattr(MSO_THEME_COLOR, theme_name, None)
        }
    return {"type": "unknown"}


def apply_fill_color(shape_obj: Shape, color_str: Optional[str]):
    if color_str in [None, "None"]:
        return
    color_info = parse_color(color_str)
    if color_info["type"] == "rgb":
        shape_obj.fill.solid()
        shape_obj.fill.fore_color.rgb = color_info["value"]
    elif color_info["type"] == "theme":
        shape_obj.fill.solid()
        shape_obj.fill.fore_color.theme_color = color_info["value"]


def apply_border(shape_obj: Shape, border: Optional[DL_Border]):
    if not hasattr(shape_obj, "line") or not border:
        return

    line = shape_obj.line
    color_info = parse_color(border.color)
    if color_info["type"] == "rgb":
        line.fill.solid()
        line.color.rgb = color_info["value"]
    elif color_info["type"] == "theme":
        line.fill.solid()
        line.color.theme_color = color_info["value"]
    elif color_info["type"] == "none":
        line.fill.background()

    if border.width_pt not in [None, "Default"]:
        try:
            line.width = Pt(float(border.width_pt))
        except:
            pass


def apply_run(run: _Run, run_data: DL_Run):
    run.text = run_data.text or ""
    font = run.font
    if run_data.font_name:
        font.name = run_data.font_name
        if font._element is not None:
            font._element.set("typeface", run_data.font_name)
    if run_data.font_size:
        font.size = Pt(run_data.font_size)
    font.bold = run_data.bold
    font.italic = run_data.italic

    color_info = parse_color(run_data.font_color)
    if color_info["type"] == "rgb":
        font.color.rgb = color_info["value"]
    elif color_info["type"] == "theme":
        font.color.theme_color = color_info["value"]


def apply_paragraph(p: _Paragraph, para: DL_TextParagraph):
    from pptx.util import Pt
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.oxml.ns import qn

    # 1. Căn lề
    p.alignment = PP_ALIGN(para.alignment or PP_ALIGN.LEFT)

    # 2. Level (bullet cấp mấy)
    if para.level is not None:
        p.level = para.level
    elif para.bullet is not None:
        p.level = para.bullet

    # 4. Line spacing (khoảng cách giữa các dòng)
    if para.line_spacing is not None:
        p.line_spacing = para.line_spacing  # là float, ví dụ 1.15 = 115%

    # 5. Bullet (ký tự hoặc kiểu đánh số)
    pPr = p._element.get_or_add_pPr()

    pPr = p._element.get_or_add_pPr()

    # Set indent qua XML attribute
    if para.left_indent is not None:
        emu_val = int(para.left_indent * 12700)
        pPr.set("marL", str(emu_val))

    if para.first_line_indent is not None:
        emu_val = int(para.first_line_indent * 12700)
        pPr.set("indent", str(emu_val))

    # Xoá bullet cũ nếu có
    for tag in ["a:buChar", "a:buAutoNum"]:
        old = pPr.find(qn(tag))
        if old is not None:
            pPr.remove(old)

    if para.bullet_type == "char" and para.bullet_char:
        buChar = OxmlElement("a:buChar")
        buChar.set("char", para.bullet_char)
        pPr.append(buChar)

    elif para.bullet_type == "number" and para.number_type:
        buAutoNum = OxmlElement("a:buAutoNum")
        buAutoNum.set("type", para.number_type)
        pPr.append(buAutoNum)

    # 6. Các run bên trong đoạn
    for run_data in para.runs:
        run = p.add_run()
        apply_run(run, run_data)


def apply_text_detail(text_frame: TextFrame, text: DL_Text):
    # Khôi phục thông tin format của text frame
    if text.frame_format:
        fmt = text.frame_format

        if fmt.wrap is not None:
            text_frame.word_wrap = fmt.wrap

        if fmt.auto_fit is not None:
            # Hiện tại python-pptx chưa hỗ trợ set auto_size trực tiếp
            # Có thể cần skip hoặc handle qua XML nếu cần sau
            pass

        if fmt.vertical_anchor:
            try:
                from pptx.enum.text import MSO_VERTICAL_ANCHOR
                text_frame.vertical_anchor = fmt.vertical_anchor
            except (KeyError, ValueError):
                pass  # fallback nếu giá trị không hợp lệ

        if fmt.margin:
            if "left" in fmt.margin:
                text_frame.margin_left = fmt.margin["left"]
            if "right" in fmt.margin:
                text_frame.margin_right = fmt.margin["right"]
            if "top" in fmt.margin:
                text_frame.margin_top = fmt.margin["top"]
            if "bottom" in fmt.margin:
                text_frame.margin_bottom = fmt.margin["bottom"]

    # Xoá nội dung cũ và khôi phục các đoạn văn bản
    text_frame.clear()

    for para in text.paragraphs:
        if para is None:
            continue
        if len(text_frame.paragraphs) == 0:
            p = text_frame.add_paragraph()
        elif len(text_frame.paragraphs) == 1 and len(text_frame.paragraphs[0].runs) == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        apply_paragraph(p, para)


def rebuild_table(shape_data: DL_Shape, slide: Slide):
    tbl_info = shape_data.table
    rows, cols = tbl_info.rows, tbl_info.cols
    pos = shape_data.position
    shape = slide.shapes.add_table(
        rows, cols, pos.x, pos.y, pos.width, pos.height)
    tbl = shape.table

    # Set col widths
    for c in range(min(cols, len(tbl_info.col_widths))):
        tbl.columns[c].width = tbl_info.col_widths[c]

    # Set row heights
    for r in range(min(rows, len(tbl_info.row_heights))):
        tbl.rows[r].height = tbl_info.row_heights[r]

    shape.width = pos.width
    shape.height = pos.height

    # Merge cells
    merged_cells = set()
    for merge in tbl_info.merge_info:
        r, c = merge.row, merge.col
        row_span = merge.row_span
        col_span = merge.col_span
        target = tbl.cell(r + row_span - 1, c + col_span - 1)
        tbl.cell(r, c).merge(target)
        for i in range(r, r + row_span):
            for j in range(c, c + col_span):
                if (i, j) != (r, c):
                    merged_cells.add((i, j))

    for r in range(rows):
        for c in range(cols):
            if (r, c) in merged_cells:
                continue
            cell = tbl.cell(r, c)
            cell.text_frame.word_wrap = True

            # Border
            if tbl_info.cell_borders and tbl_info.cell_borders[r][c]:
                apply_cell_border(cell, tbl_info.cell_borders[r][c])

            # Text
            if tbl_info.data_detail and tbl_info.data_detail[r][c]:
                apply_text_detail(cell.text_frame, tbl_info.data_detail[r][c])

            # Fill
            if tbl_info.cell_fills and tbl_info.cell_fills[r][c]:
                apply_fill_color(cell, tbl_info.cell_fills[r][c])

    return shape


def rebuild_textbox(shape_data: DL_Shape, slide: Slide):
    pos = shape_data.position
    shape_type = shape_data.type
    if shape_type == 1:
        autoshape_type_id = 1
        shape = slide.shapes.add_shape(
            autoshape_type_id, pos.x, pos.y, pos.width, pos.height)
    else:
        shape = slide.shapes.add_textbox(pos.x, pos.y, pos.width, pos.height)
    tf = shape.text_frame
    if shape_data.text:
        apply_text_detail(tf, shape_data.text)
    return shape


def rebuild_image(shape_data: Picture, slide: Slide, json_path: str):
    if not shape_data.image or not shape_data.image.filename:
        raise ValueError(
            f"[Slide {shape_data.shape_index}] Thiếu thông tin image để khôi phục")

    # Tính full path từ đường dẫn tương đối trong JSON
    json_folder = os.path.dirname(json_path)
    image_path = os.path.join(json_folder, shape_data.image.filename)

    if not os.path.isfile(image_path):
        raise FileNotFoundError(f"Không tìm thấy file ảnh: {image_path}")

    pos = shape_data.position
    pic = slide.shapes.add_picture(
        image_path,
        pos.x, pos.y,
        width=pos.width,
        height=pos.height
    )
    return pic


def build_pptx_from_json(json_path: str, output_path: str):
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    pptx_data = from_dict(data_class=DL_PPTXData, data=data)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    prs.slide_width = pptx_data.slide_width
    prs.slide_height = pptx_data.slide_height

    for slide_data in pptx_data.slides:
        slide = prs.slides.add_slide(blank_layout)
        for shape_data in slide_data.shapes:
            shape = None
            if shape_data.table:
                shape = rebuild_table(shape_data, slide)
            elif shape_data.text:
                shape = rebuild_textbox(shape_data, slide)
            elif shape_data.image:
                shape = rebuild_image(shape_data, slide, json_path)
            if shape:
                apply_fill_color(shape, shape_data.background_fill_color)
                apply_border(shape, shape_data.border)

    prs.save(output_path)
    print(f"✅ PPTX đã được tạo tại: {output_path}")


if __name__ == "__main__":
    # build_pptx_from_json(r"utest\dump_test_ppt1.json",
    #                      "bin/test_ppt1_restored_from_json.pptx")
    # build_pptx_from_json(r"bin\Pre_DOI_Form_05_2024v2\Pre_DOI_Form_05_2024v2.json",
    #                      r"bin\Pre_DOI_Form_05_2024v2\Pre_DOI_Form_05_2024v2.json.pptx")
    build_pptx_from_json(r"bin\Pre_DOI_Form_05_2024v3\Pre_DOI_Form_05_2024v3.json",
                         r"bin\Pre_DOI_Form_05_2024v3\Pre_DOI_Form_05_2024v3.json.pptx")