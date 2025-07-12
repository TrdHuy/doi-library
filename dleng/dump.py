# Refactored version of dump.py focusing on modularization
import hashlib
import json
import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_VERTICAL_ANCHOR
from pptx.oxml.ns import qn
from pptx.dml.fill import _NoFill
from pptx.dml.fill import _NoneFill
from pptx.shapes.picture import Picture
from utils.pptxhelper import *

SHAPE_TYPES_WITH_FILL_LINE = {
    MSO_SHAPE_TYPE.AUTO_SHAPE,
    MSO_SHAPE_TYPE.FORM_CONTROL,
    MSO_SHAPE_TYPE.PICTURE,
    MSO_SHAPE_TYPE.PLACEHOLDER,
    MSO_SHAPE_TYPE.TEXT_BOX,
    MSO_SHAPE_TYPE.FREEFORM
    }
TAG_MAP = {
    "left": "lnL",
    "right": "lnR",
    "top": "lnT",
    "bottom": "lnB",
    "diagonal_down": "lnTlToBr",
    "diagonal_up": "lnBlToTr"
}


def safe_deep_dump(obj, max_depth=2, _visited=None, _depth=0):
    if _visited is None:
        _visited = set()
    if id(obj) in _visited:
        return "[Circular Reference]"
    _visited.add(id(obj))
    if isinstance(obj, (str, int, float, bool, type(None))):
        return obj
    if _depth >= max_depth:
        return f"<{type(obj).__name__}>"
    result = {}
    for attr in dir(obj):
        if attr.startswith("_"):
            continue
        try:
            value = getattr(obj, attr)
            if callable(value):
                continue
            result[attr] = safe_deep_dump(
                value, max_depth, _visited, _depth + 1)
        except Exception as e:
            result[attr] = f"[Error: {e}]"
    return result


def get_rgb_safe(color_obj, context="(unknown context)"):
    if color_obj is None:
        return "None"
    try:
        if color_obj.rgb:
            return f"RGB:{color_obj.rgb}"
    except AttributeError:
        pass
    try:
        if color_obj.theme_color:
            theme_name = color_obj.theme_color.name
            raise ValueError(
                f"{context} – sử dụng màu theo theme ({theme_name}) mà không có RGB cụ thể.")
    except:
        pass
    raise ValueError(f"{context} – không xác định được màu fill/font.")


def extract_shape_border_info(shape, context):
    if not hasattr(shape, "line") or shape.line is None:
        raise ValueError(f"{context} không có shape.line")
    line = shape.line
    line_fill = line.fill
    no_outline = line_fill is None or line_fill.type is None or isinstance(
        line.fill._fill, _NoFill)
    if not no_outline:
        try:
            if line_fill.fore_color is None or line_fill.fore_color.rgb is None:
                no_outline = True
        except AttributeError:
            no_outline = True
    if no_outline:
        return {"color": "None", "width_pt": "Default", "style": "None"}
    return {
        "color": get_rgb_safe(line.color, context=f"{context} border.color"),
        "width_pt": round(line.width / 12700, 2) if line.width else "Default",
        "style": str(line.dash_style) if line.dash_style else "None"
    }


def extract_run_info(run, context):
    font = run.font
    font_size_pt = font.size.pt if font.size else None
    if font_size_pt is None:
        raise ValueError(f"{context} thiếu font size rõ ràng")
    font_name = font.name or (font._element.get(
        "typeface") if font._element is not None else None)
    if font_name is None:
        raise ValueError(f"{context} thiếu font name rõ ràng")
    return {
        "text": run.text,
        "font_name": font_name,
        "font_size": font_size_pt,
        "bold": font.bold,
        "italic": font.italic,
        "font_color": get_rgb_safe(font.color, context=context)
    }


def extract_paragraph_info(paragraph, context):
    alignment_val = paragraph.alignment or PP_ALIGN.LEFT
    font = paragraph.font
    
    font_size_pt = font.size.pt if font.size else None
    font_name = font.name if font.name else None
    bold = font.bold if font.bold is not None else None
    italic = font.italic if font.italic is not None else None
    font_color = get_rgb_safe(font.color, context=context) if font.color.type else None

    # Nếu thiếu thông tin nào → fallback bằng endParaRPr
    if not all([font_size_pt, font_name, bold is not None, italic is not None, font_color]):
        if paragraph.runs:
            first_run_font = paragraph.runs[0].font
            font_size_pt = first_run_font.size.pt if first_run_font.size else None
            font_name = first_run_font.name if first_run_font.name else None
            bold = first_run_font.bold if first_run_font.bold is not None else None
            italic = first_run_font.italic if first_run_font.italic is not None else None
            font_color = get_rgb_safe(first_run_font.color, context=context) if first_run_font.color else None
        else:
            fallback = extract_font_info_from_end_para(paragraph)
            if fallback:
                font_size_pt = font_size_pt or fallback["font_size"]
                font_name = font_name or fallback["font_name"]
                bold = bold if bold is not None else fallback["bold"]
                italic = italic if italic is not None else fallback["italic"]
                font_color = font_color or fallback["font_color"]

    if font_size_pt is None:
        raise ValueError(f"{context} thiếu font size rõ ràng")

    para_info = {
        "alignment": alignment_val,
        "runs": [],
        "bullet": paragraph.level,
        "bullet_type": None,
        "font_name": font_name,
        "font_size": font_size_pt,
        "bold": bold,
        "italic": italic,
        "font_color": font_color
    }

    if paragraph.level is not None:
        if paragraph._pPr is not None:
            buChar = paragraph._pPr.find(qn("a:buChar"))
            buAutoNum = paragraph._pPr.find(qn("a:buAutoNum"))
            if buChar is not None:
                para_info["bullet_type"] = "char"
                para_info["bullet_char"] = buChar.attrib.get("char", "")
            elif buAutoNum is not None:
                para_info["bullet_type"] = "number"
                para_info["number_type"] = buAutoNum.attrib.get("type", "arabicPeriod")
    
    if paragraph._pPr is not None:
        marL = paragraph._pPr.attrib.get("marL")
        indent = paragraph._pPr.attrib.get("indent")
        level = paragraph._pPr.attrib.get("lvl")
        para_info["level"] = int(level) if level is not None else 0
        para_info["left_indent"] = round(int(marL) / 12700, 2) if marL else None
        para_info["first_line_indent"] = round(int(indent) / 12700, 2) if indent else None
    else:
        para_info["left_indent"] = None
        para_info["first_line_indent"] = None
        para_info["level"] = None

    lnSpc = paragraph._pPr.find(qn("a:lnSpc"))
    if lnSpc is not None:
        spcPct = lnSpc.find(qn("a:spcPct"))
        if spcPct is not None and "val" in spcPct.attrib:
            para_info["line_spacing"] = int(spcPct.attrib["val"]) / 100000  # chuyển về tỉ lệ thực
        else:
            para_info["line_spacing"] = None
    else:
        para_info["line_spacing"] = None

    for run_idx, run in enumerate(paragraph.runs):
        run_ctx = f"{context} - Run {run_idx+1}"
        run_info = extract_run_info(run, run_ctx)
        run_info["run_index"] = run_idx + 1
        para_info["runs"].append(run_info)
    return para_info


def extract_cell_text_detail(cell, slide_idx, shape_idx, r_idx, c_idx):
    context = f"[Slide {slide_idx+1} - Shape {shape_idx+1} - Cell ({r_idx+1},{c_idx+1})]"
    detail = []
    for p_idx, para in enumerate(cell.text_frame.paragraphs):
        para_ctx = f"{context} - Para {p_idx+1}"
        para_info = extract_paragraph_info(para, para_ctx)
        para_info["paragraph_index"] = p_idx + 1
        detail.append(para_info)
    return {
        "frame_format": extract_text_frame_format(cell.text_frame),
        "paragraphs": detail
    }


def extract_cell_border(cell, slide_idx, shape_idx, r_idx, c_idx):
    tcPr = cell._tc.tcPr
    borders = {}
    for side, tag in TAG_MAP.items():
        ln = tcPr.find(qn(f'a:{tag}'))
        border_info = {"color": "None",
                       "width": "Default", "dash_type": "None"}
        if ln is not None:
            solid_fill = ln.find(qn("a:solidFill"))
            if solid_fill is not None:
                srgb = solid_fill.find(qn("a:srgbClr"))
                if srgb is not None:
                    hex_val = srgb.attrib.get("val")
                    border_info["color"] = f"RGB:{hex_val.upper()}"
            if "w" in ln.attrib:
                try:
                    border_info["width"] = round(
                        int(ln.attrib["w"]) / 12700, 2)
                except:
                    border_info["width"] = "Error"
            prst_dash = ln.find(qn("a:prstDash"))
            if prst_dash is not None:
                border_info["dash_type"] = prst_dash.attrib.get("val", "None")
        borders[side] = border_info
    return borders


def extract_text_frame_format(text_frame):
    format_info = {}

    # Có wrap word không?
    format_info["wrap"] = text_frame.word_wrap

    # Auto-fit text?
    format_info["auto_fit"] = (
        text_frame.auto_size is not None  # can be `None`, `TextAutoSize.SHAPE_TO_FIT_TEXT`, ...
    )

    # Vertical alignment
    if text_frame.vertical_anchor:
        format_info["vertical_anchor"] = int(text_frame.vertical_anchor)
    else:
        format_info["vertical_anchor"] = int(MSO_VERTICAL_ANCHOR.TOP)

    # Margins (EMU)
    format_info["margin"] = {
        "left": text_frame.margin_left,
        "right": text_frame.margin_right,
        "top": text_frame.margin_top,
        "bottom": text_frame.margin_bottom
    }

    return format_info


def extract_text_from_shape(shape, slide_idx, shape_idx, for_txt):
    tf = shape.text_frame
    paragraphs = []
    for p_idx, para in enumerate(tf.paragraphs):
        context = f"[Slide {slide_idx+1} - Shape {shape_idx+1} - Paragraph {p_idx+1}]"
        para_info = extract_paragraph_info(para, context)
        para_info["paragraph_index"] = p_idx + 1
        para_info["text"] = para.text.strip().replace("\n", "\\n") if for_txt else para.text.strip()
        paragraphs.append(para_info)

    return {
        "frame_format": extract_text_frame_format(tf),
        "paragraphs": paragraphs
    }

def extract_table_from_shape(shape, slide_idx, shape_idx, for_txt):
    tbl = shape.table
    num_rows = len(tbl.rows)
    num_cols = len(tbl.columns)
    table_data = [["" for _ in range(num_cols)] for _ in range(num_rows)]
    table_data_detail = [
        [None for _ in range(num_cols)] for _ in range(num_rows)]
    cell_fills = [["None" for _ in range(num_cols)] for _ in range(num_rows)]
    merge_info = []
    cell_borders = [[None for _ in range(num_cols)] for _ in range(num_rows)]

    for r_idx in range(num_rows):
        for c_idx in range(num_cols):
            cell = tbl.cell(r_idx, c_idx)
            if cell.is_spanned and not cell.is_merge_origin:
                continue
            table_data[r_idx][c_idx] = cell.text.strip().replace(
                "\n", "\\n") if for_txt else cell.text.strip()
            table_data_detail[r_idx][c_idx] = extract_cell_text_detail(
                cell, slide_idx, shape_idx, r_idx, c_idx)
            if not hasattr(cell, "fill") or cell.fill is None or isinstance(cell.fill._fill, _NoneFill):
                raise ValueError(
                    f"[Slide {slide_idx+1} - Shape {shape_idx+1} - Cell ({r_idx+1},{c_idx+1})] thiếu fill")

            if not hasattr(cell, "fill") or cell.fill is None or not cell.fill.fore_color:
                raise ValueError(
                    f"[Slide {slide_idx+1} - Shape {shape_idx+1} - Cell ({r_idx+1},{c_idx+1})] thiếu fill")
            cell_fills[r_idx][c_idx] = get_rgb_safe(
                cell.fill.fore_color, context=f"[Slide {slide_idx+1} - Shape {shape_idx+1} - Cell ({r_idx+1},{c_idx+1})]")
            if cell.span_height > 1 or cell.span_width > 1:
                merge_info.append(
                    {"row": r_idx, "col": c_idx, "row_span": cell.span_height, "col_span": cell.span_width})
            cell_borders[r_idx][c_idx] = extract_cell_border(
                cell, slide_idx, shape_idx, r_idx, c_idx)

    col_widths = [col.width for col in tbl.columns]
    row_heights = [row.height for row in tbl.rows]

    return {
        "rows": num_rows,
        "cols": num_cols,
        "data": table_data,
        "data_detail": table_data_detail,
        "cell_fills": cell_fills,
        "merge_info": merge_info,
        "col_widths": col_widths,
        "row_heights": row_heights,
        "cell_borders": cell_borders
    }


def extract_picture_info(shape: Picture, slide_idx, shape_idx, asset_dir):
    image = shape.image
    ext = image.ext.strip(".")
    img_bytes = image.blob

    if not img_bytes:
        raise ValueError(
            f"[Slide {slide_idx+1} - Shape {shape_idx+1}] – Không có dữ liệu ảnh")

    hash_part = hashlib.md5(img_bytes).hexdigest()[:8]
    export_name = f"img_slide{slide_idx+1}_shape{shape_idx+1}_{hash_part}.{ext}"
    export_path = os.path.join(asset_dir, export_name)

    with open(export_path, "wb") as f:
        f.write(img_bytes)

    return {
        "filename": os.path.join("asset", export_name),
        "ext": ext,
        "content_type": image.content_type,
        "size": len(img_bytes)
    }


def extract_slide_data(pptx_path, output_dir, for_txt=False, is_debug=False):
    prs = Presentation(pptx_path)
    slides = []

    asset_dir = os.path.join(output_dir, "asset")
    os.makedirs(asset_dir, exist_ok=True)

    for i, slide in enumerate(prs.slides):
        slide_info = {"slide_number": i + 1, 
                      "shapes": [],
                      "slide_id": str(slide.slide_id)}
        
        for j, shape in enumerate(slide.shapes):
            shape_info = {
                "shape_name": shape.name,
                "shape_index": j + 1,
                "type": shape.shape_type,
                "position": {
                    "x": shape.left,
                    "y": shape.top,
                    "width": shape.width,
                    "height": shape.height
                },
                "background_fill_color": None,
                "border": {},
                "text": None,
                "table": None,
                "shape_visible": get_shape_visible(shape)
            }

            if is_debug:
                shape_info["raw_attributes"] = safe_deep_dump(
                    shape, max_depth=2)

            shape_type = shape.shape_type

            if shape_type == MSO_SHAPE_TYPE.GROUP:
                raise ValueError(
                    f"[Slide {slide_info}] – Không hỗ trợ dump cho shape kiểu group, vui lòng bỏ group")

            has_visual_style = shape_type in SHAPE_TYPES_WITH_FILL_LINE

            if has_visual_style and hasattr(shape, "fill") and shape.fill and shape.fill.fore_color:
                shape_info["background_fill_color"] = get_rgb_safe(
                    shape.fill.fore_color, context=f"[Slide {i+1} - Shape {j+1}] fill")
                shape_info["border"] = extract_shape_border_info(
                    shape, context=f"[Slide {i+1} - Shape {j+1}]")
            
            # Nếu là shape SLIDE_INFO → parse json trong text frame
            if shape.name == "SLIDE_INFO" and shape.has_text_frame:
                try:
                    slide_tag_json = shape.text_frame.text.strip()
                    slide_info["slide_tag_info"] = json.loads(slide_tag_json)
                except Exception as e:
                    print(f"[Slide {i+1}] ❌ Không parse được JSON từ SLIDE_INFO: {e}")
                    slide_info["slide_tag_info"] = None
            elif shape.has_text_frame:
                shape_info["text"] = extract_text_from_shape(
                    shape, i, j, for_txt)

            if shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_info["table"] = extract_table_from_shape(
                    shape, i, j, for_txt)

            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_info["image"] = extract_picture_info(
                    shape, i, j, asset_dir)

            slide_info["shapes"].append(shape_info)
        slides.append(slide_info)

    return {
        "slide_width": prs.slide_width,
        "slide_height": prs.slide_height,
        "slides": slides
    }


def describe_pptx_to_json_with_assets(pptx_path, output_root_folder):
    slide_name = os.path.splitext(os.path.basename(pptx_path))[0]
    output_dir = os.path.join(output_root_folder, slide_name)
    os.makedirs(output_dir, exist_ok=True)
    json_path = os.path.join(output_dir, f"{slide_name}.json")

    data = extract_slide_data(pptx_path, output_dir)
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2, ensure_ascii=False)


# Ví dụ sử dụng
if __name__ == "__main__":

    from pathlib import Path
    # Đường dẫn tuyệt đối của file đang chạy
    current_file = Path(__file__).resolve()
    # Lấy root folder (ở đây là cha của thư mục 'dleng')
    root_dir = current_file.parents[1]
    # Ví dụ sử dụng
    describe_pptx_to_json_with_assets(
        pptx_path=root_dir / "template" / "Pre_DOI_Form_05_2024_v3.pptx",
        output_root_folder=root_dir / "template")
    # describe_pptx_to_json_with_assets(
    #     r"dleng\utest\test_ppt1.pptx", "bin")
    # describe_pptx_to_json(
    #     "utest\\test_ppt1.pptx", "bin\\test_ppt1.json")
    # describe_pptx_to_json("bin\\test_ppt1_restored_from_json.pptx",
    #                       "bin\\dump_test_ppt1_restored_from_json.json")
