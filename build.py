# Refactored build.py with modular table and text rebuild logic

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.oxml.xmlchemy import OxmlElement

EMU = 1  # đơn vị đã là EMU trong JSON dump


def apply_shape_format(shape_obj, raw_attrs):
    # Fill
    try:
        fill_rgb = raw_attrs.get("fill", {}).get("fore_color")
        if fill_rgb and "RGBColor" in str(fill_rgb):
            r, g, b = map(int, fill_rgb.strip("RGBColor() ").split(","))
            shape_obj.fill.solid()
            shape_obj.fill.fore_color.rgb = RGBColor(r, g, b)
    except:
        pass

    # Border
    try:
        line_info = raw_attrs.get("line", {})
        line_rgb = line_info.get("color")
        if line_rgb and "RGBColor" in str(line_rgb):
            r, g, b = map(int, line_rgb.strip("RGBColor() ").split(","))
            shape_obj.line.color.rgb = RGBColor(r, g, b)
        width = line_info.get("width")
        if width and isinstance(width, (int, float)):
            shape_obj.line.width = Pt(width)
    except:
        pass

    # TextFrame
    try:
        tf_attrs = raw_attrs.get("text_frame", {})
        tf = shape_obj.text_frame
        tf.margin_top = tf_attrs.get("margin_top", tf.margin_top)
        tf.margin_bottom = tf_attrs.get("margin_bottom", tf.margin_bottom)
        tf.margin_left = tf_attrs.get("margin_left", tf.margin_left)
        tf.margin_right = tf_attrs.get("margin_right", tf.margin_right)
        tf.auto_size = tf_attrs.get("auto_size", tf.auto_size)
        tf.word_wrap = tf_attrs.get("word_wrap", tf.word_wrap)
    except:
        pass


def apply_cell_border(cell, border_info):
    tcPr = cell._tc.get_or_add_tcPr()
    side_map = {
        "left": "a:lnL",
        "right": "a:lnR",
        "top": "a:lnT",
        "bottom": "a:lnB",
        "tl2br": "a:lnTlToBr",
        "bl2tr": "a:lnBlToTr"
    }

    for side, tag in side_map.items():
        if side not in border_info:
            continue

        side_border = border_info[side]
        color_str = side_border.get("color")
        width = side_border.get("width")
        dash_style = side_border.get("dash_style")

        ln = OxmlElement(tag)
        if width and isinstance(width, (int, float)):
            ln.set("w", str(int(width * 12700)))  # convert pt to EMU

        # Dash style
        if dash_style and dash_style.lower() != "none":
            prstDash = OxmlElement("a:prstDash")
            prstDash.set("val", dash_style)
            ln.append(prstDash)

        # Color
        color_info = parse_color(color_str)
        if color_info["type"] == "rgb":
            solidFill = OxmlElement("a:solidFill")
            srgbClr = OxmlElement("a:srgbClr")
            rgb_val = color_info["value"]
            srgbClr.set(
                "val", f"{rgb_val[0]:02X}{rgb_val[1]:02X}{rgb_val[2]:02X}")
            solidFill.append(srgbClr)
            ln.append(solidFill)
        elif color_info["type"] == "theme":
            solidFill = OxmlElement("a:solidFill")
            schemeClr = OxmlElement("a:schemeClr")
            schemeClr.set("val", color_info["value"].name.lower())
            solidFill.append(schemeClr)
            ln.append(solidFill)

        tcPr.append(ln)


def parse_color(color_str):
    if color_str is None or color_str == "None":
        return {"type": "none"}
    if color_str.startswith("RGB:"):
        hex_str = color_str[4:].strip()
        return {
            "type": "rgb",
            "value": RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))
        }
    if color_str.startswith("Theme:"):
        theme_name = color_str[6:].strip()
        return {
            "type": "theme",
            "value": getattr(MSO_THEME_COLOR, theme_name, None)
        }
    return {"type": "unknown"}


def apply_fill_color(obj, color_str):
    if color_str in [None, "None"]:
        return
    color_info = parse_color(color_str)
    if color_info["type"] == "rgb":
        obj.fill.solid()
        obj.fill.fore_color.rgb = color_info["value"]
    elif color_info["type"] == "theme":
        obj.fill.solid()
        obj.fill.fore_color.theme_color = color_info["value"]


def apply_border(shape_obj, border_dict):
    if not hasattr(shape_obj, "line") or not border_dict:
        return
    line = shape_obj.line
    color_info = parse_color(border_dict.get("color"))
    if color_info["type"] == "rgb":
        line.fill.solid()
        line.color.rgb = color_info["value"]
    elif color_info["type"] == "theme":
        line.fill.solid()
        line.color.theme_color = color_info["value"]
    elif color_info["type"] == "none":
        line.fill.background()

    if border_dict.get("width_pt") not in [None, "Default"]:
        try:
            line.width = Pt(float(border_dict.get("width_pt")))
        except:
            pass


def apply_run(run, run_data):
    run.text = run_data.get("text", "")
    font = run.font
    if run_data.get("font_name"):
        font.name = run_data["font_name"]
        if font._element is not None:
            font._element.set("typeface", run_data["font_name"])
    if run_data.get("font_size"):
        font.size = Pt(run_data["font_size"])
    font.bold = run_data.get("bold", False)
    font.italic = run_data.get("italic", False)
    color_info = parse_color(run_data.get("font_color"))
    if color_info["type"] == "rgb":
        font.color.rgb = color_info["value"]
    elif color_info["type"] == "theme":
        font.color.theme_color = color_info["value"]


def apply_paragraph(p, para_data):
    p.alignment = PP_ALIGN(para_data.get("alignment", PP_ALIGN.LEFT))
    p.font.bold = False
    p.font.italic = False
    for run_data in para_data.get("runs", []):
        run = p.add_run()
        apply_run(run, run_data)


def apply_text_detail(text_frame, detail):
    text_frame.clear()
    for para_data in detail:
        if para_data is None:
            continue
        if len(text_frame.paragraphs) == 0:
            p = text_frame.add_paragraph()
        elif len(text_frame.paragraphs) == 1 and len(text_frame.paragraphs[0].runs) == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        apply_paragraph(p, para_data)


def rebuild_table(shape_data, slide):
    tbl_info = shape_data["table"]
    rows, cols = tbl_info["rows"], tbl_info["cols"]
    x, y, w, h = shape_data["position"].values()
    shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = shape.table

    # Set col widths
    col_widths = tbl_info.get("col_widths", [])
    for c in range(min(cols, len(col_widths))):
        tbl.columns[c].width = col_widths[c]

    # Set row heights
    row_heights = tbl_info.get("row_heights", [])
    for r in range(min(rows, len(row_heights))):
        tbl.rows[r].height = row_heights[r]
    shape.width = w
    shape.height = h
    
    # Merge cells
    merged_cells = set()
    for merge in tbl_info.get("merge_info", []):
        r, c = merge["row"], merge["col"]
        row_span = merge.get("row_span", 1)
        col_span = merge.get("col_span", 1)
        target = tbl.cell(r + row_span - 1, c + col_span - 1)
        tbl.cell(r, c).merge(target)
        for i in range(r, r + row_span):
            for j in range(c, c + col_span):
                if (i, j) != (r, c):
                    merged_cells.add((i, j))

    # Fill cell data
    for r in range(rows):
        for c in range(cols):
            if (r, c) in merged_cells:
                continue
            cell = tbl.cell(r, c)
            cell.text_frame.word_wrap = True

            # Border
            border_info = tbl_info.get(
                "cell_borders", [[None]*cols]*rows)[r][c]
            if border_info:
                apply_cell_border(cell, border_info)

            # Text
            detail = tbl_info.get("data_detail", [[[]]*cols]*rows)
            if r < len(detail) and c < len(detail[r]):
                apply_text_detail(cell.text_frame, detail[r][c])

            # Fill
            fills = tbl_info.get("cell_fills", [["None"]*cols]*rows)
            if r < len(fills) and c < len(fills[r]):
                apply_fill_color(cell, fills[r][c])

    return shape


def rebuild_textbox(shape_data, slide):
    x, y, w, h = shape_data["position"].values()
    shape_type = shape_data.get("type", -1)
    if shape_type == 1:
        shape = slide.shapes.add_shape(
            autoshape_type_id=shape_data.get(
                "raw_attributes", {}).get("auto_shape_type", 1),
            left=x, top=y, width=w, height=h
        )
    else:
        shape = slide.shapes.add_textbox(x, y, w, h)
    tf = shape.text_frame
    apply_text_detail(tf, shape_data.get("text", []))
    return shape


def build_pptx_from_json(json_path, output_path):
    with open(json_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    prs.slide_width = data["slide_width"]
    prs.slide_height = data["slide_height"]

    for slide_data in data["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        for shape_data in slide_data["shapes"]:
            if "table" in shape_data and shape_data["table"]:
                shape = rebuild_table(shape_data, slide)
            elif "text" in shape_data and shape_data["text"]:
                shape = rebuild_textbox(shape_data, slide)
            else:
                continue
            #apply_shape_format(shape, shape_data.get("raw_attributes", {}))
            apply_fill_color(shape, shape_data.get("background_fill_color"))
            apply_border(shape, shape_data.get("border"))

    prs.save(output_path)
    print(f"✅ PPTX đã được tạo tại: {output_path}")


if __name__ == "__main__":
    build_pptx_from_json("bin/test_ppt1.json", "bin/test_ppt1_restored_from_json.pptx")
